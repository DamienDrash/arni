"""Multi-format document parser registry.

Implements the Parser Registry pattern: each file format is handled by
a dedicated parser class.  The registry auto-discovers parsers and
dispatches files to the correct one based on extension and MIME type.

Supported formats (Phase 1):
    .md, .txt       – Plain text / Markdown
    .pdf            – PDF documents
    .docx           – Microsoft Word
    .html, .htm     – HTML pages

Supported formats (Phase 3 – extended):
    .xlsx, .xls     – Excel spreadsheets
    .csv, .tsv      – Delimited data
    .pptx           – PowerPoint presentations
    .eml            – Email messages
    .rtf            – Rich Text Format
    .json           – JSON data
"""

from __future__ import annotations

import os
from abc import ABC, abstractmethod
from typing import Any

import structlog

from app.memory_platform.models import ContentChunk

logger = structlog.get_logger()


class BaseParser(ABC):
    """Abstract base class for all document parsers."""

    # Subclasses MUST set these
    supported_extensions: list[str] = []
    supported_mimetypes: list[str] = []
    parser_name: str = "base"

    @abstractmethod
    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        """Parse the file at *file_path* and return a list of content chunks."""
        ...

    def can_handle(self, filename: str, mimetype: str = "") -> bool:
        """Check whether this parser can handle the given file."""
        ext = os.path.splitext(filename)[1].lower()
        if ext in self.supported_extensions:
            return True
        if mimetype and mimetype in self.supported_mimetypes:
            return True
        return False


class MarkdownParser(BaseParser):
    """Parser for Markdown and plain text files."""

    supported_extensions = [".md", ".txt", ".text"]
    supported_mimetypes = ["text/markdown", "text/plain"]
    parser_name = "markdown"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        except Exception as exc:
            logger.error("parser.markdown.read_error", path=file_path, error=str(exc))
            return []

        if not content.strip():
            return []

        chunks = self._split_by_sections(content, metadata)
        logger.info("parser.markdown.parsed", path=file_path, chunks=len(chunks))
        return chunks

    def _split_by_sections(
        self,
        content: str,
        metadata: dict[str, Any],
    ) -> list[ContentChunk]:
        """Split Markdown content by headings into semantic chunks."""
        lines = content.split("\n")
        chunks: list[ContentChunk] = []
        current_section = ""
        current_title = ""
        section_start = 0

        for i, line in enumerate(lines):
            if line.startswith("#"):
                if current_section.strip():
                    chunks.append(ContentChunk(
                        content=current_section.strip(),
                        content_type="text",
                        section_title=current_title or None,
                        metadata={**metadata, "line_start": section_start},
                    ))
                current_title = line.lstrip("#").strip()
                current_section = line + "\n"
                section_start = i
            else:
                current_section += line + "\n"

        # Last section
        if current_section.strip():
            chunks.append(ContentChunk(
                content=current_section.strip(),
                content_type="text",
                section_title=current_title or None,
                metadata={**metadata, "line_start": section_start},
            ))

        # If no sections found, chunk by size
        if len(chunks) <= 1 and len(content) > 2000:
            return self._chunk_by_size(content, metadata)

        return chunks

    def _chunk_by_size(
        self,
        content: str,
        metadata: dict[str, Any],
        chunk_size: int = 1000,
        overlap: int = 200,
    ) -> list[ContentChunk]:
        """Fall back to fixed-size chunking with overlap."""
        chunks: list[ContentChunk] = []
        start = 0
        while start < len(content):
            end = start + chunk_size
            chunk_text = content[start:end]
            if chunk_text.strip():
                chunks.append(ContentChunk(
                    content=chunk_text.strip(),
                    content_type="text",
                    metadata={**metadata, "char_start": start},
                ))
            start = end - overlap
        return chunks


class PDFParser(BaseParser):
    """Parser for PDF documents using pdfplumber or fallback to PyPDF2."""

    supported_extensions = [".pdf"]
    supported_mimetypes = ["application/pdf"]
    parser_name = "pdf"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        chunks: list[ContentChunk] = []

        # Try pdfplumber first (better table extraction)
        try:
            import pdfplumber  # type: ignore[import-untyped]

            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        chunks.append(ContentChunk(
                            content=text.strip(),
                            content_type="text",
                            page_number=page_num,
                            metadata={**metadata, "parser": "pdfplumber"},
                        ))

                    # Extract tables
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables):
                        if table:
                            table_md = self._table_to_markdown(table)
                            if table_md.strip():
                                chunks.append(ContentChunk(
                                    content=table_md,
                                    content_type="table",
                                    page_number=page_num,
                                    metadata={
                                        **metadata,
                                        "parser": "pdfplumber",
                                        "table_index": table_idx,
                                    },
                                ))

            logger.info("parser.pdf.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except ImportError:
            pass

        # Fallback to PyPDF2
        try:
            from PyPDF2 import PdfReader  # type: ignore[import-untyped]

            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    chunks.append(ContentChunk(
                        content=text.strip(),
                        content_type="text",
                        page_number=page_num,
                        metadata={**metadata, "parser": "pypdf2"},
                    ))

            logger.info("parser.pdf.parsed_pypdf2", path=file_path, chunks=len(chunks))
            return chunks

        except ImportError:
            logger.error("parser.pdf.no_library", path=file_path)
            return []

    @staticmethod
    def _table_to_markdown(table: list[list[str | None]]) -> str:
        """Convert a table (list of rows) to Markdown format."""
        if not table or not table[0]:
            return ""
        rows = []
        for row in table:
            cells = [str(c or "").strip() for c in row]
            rows.append("| " + " | ".join(cells) + " |")
        if len(rows) > 1:
            header_sep = "| " + " | ".join(["---"] * len(table[0])) + " |"
            rows.insert(1, header_sep)
        return "\n".join(rows)


class DocxParser(BaseParser):
    """Parser for Microsoft Word (.docx) documents."""

    supported_extensions = [".docx", ".doc"]
    supported_mimetypes = [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ]
    parser_name = "docx"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        chunks: list[ContentChunk] = []

        try:
            from docx import Document  # type: ignore[import-untyped]

            doc = Document(file_path)
            current_section = ""
            current_title = ""

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detect headings
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    if current_section.strip():
                        chunks.append(ContentChunk(
                            content=current_section.strip(),
                            content_type="text",
                            section_title=current_title or None,
                            metadata={**metadata, "parser": "python-docx"},
                        ))
                    current_title = text
                    current_section = f"## {text}\n\n"
                else:
                    current_section += text + "\n\n"

            # Last section
            if current_section.strip():
                chunks.append(ContentChunk(
                    content=current_section.strip(),
                    content_type="text",
                    section_title=current_title or None,
                    metadata={**metadata, "parser": "python-docx"},
                ))

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    md = self._rows_to_markdown(rows)
                    chunks.append(ContentChunk(
                        content=md,
                        content_type="table",
                        metadata={
                            **metadata,
                            "parser": "python-docx",
                            "table_index": table_idx,
                        },
                    ))

            logger.info("parser.docx.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except ImportError:
            logger.error("parser.docx.python_docx_not_installed", path=file_path)
            return []
        except Exception as exc:
            logger.error("parser.docx.error", path=file_path, error=str(exc))
            return []

    @staticmethod
    def _rows_to_markdown(rows: list[list[str]]) -> str:
        if not rows:
            return ""
        lines = []
        for row in rows:
            lines.append("| " + " | ".join(row) + " |")
        if len(lines) > 1:
            lines.insert(1, "| " + " | ".join(["---"] * len(rows[0])) + " |")
        return "\n".join(lines)


class HTMLParser(BaseParser):
    """Parser for HTML files."""

    supported_extensions = [".html", ".htm"]
    supported_mimetypes = ["text/html"]
    parser_name = "html"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")

            # Remove script and style elements
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()

            text = soup.get_text(separator="\n", strip=True)
            if not text.strip():
                return []

            # Use markdown parser's chunking logic
            md_parser = MarkdownParser()
            chunks = md_parser._chunk_by_size(text, {**metadata, "parser": "html"})
            logger.info("parser.html.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except Exception as exc:
            logger.error("parser.html.error", path=file_path, error=str(exc))
            return []


class CSVParser(BaseParser):
    """Parser for CSV and TSV files."""

    supported_extensions = [".csv", ".tsv"]
    supported_mimetypes = ["text/csv", "text/tab-separated-values"]
    parser_name = "csv"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        try:
            import csv

            delimiter = "\t" if file_path.endswith(".tsv") else ","
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f, delimiter=delimiter)
                rows = list(reader)

            if not rows:
                return []

            # Convert to markdown table in batches of 50 rows
            chunks: list[ContentChunk] = []
            batch_size = 50
            header = rows[0]

            for i in range(1, len(rows), batch_size):
                batch = rows[i:i + batch_size]
                lines = ["| " + " | ".join(header) + " |"]
                lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in batch:
                    # Pad or truncate row to match header length
                    padded = row + [""] * (len(header) - len(row))
                    lines.append("| " + " | ".join(padded[:len(header)]) + " |")
                chunks.append(ContentChunk(
                    content="\n".join(lines),
                    content_type="table",
                    metadata={
                        **metadata,
                        "parser": "csv",
                        "row_start": i,
                        "row_end": min(i + batch_size, len(rows)),
                    },
                ))

            logger.info("parser.csv.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except Exception as exc:
            logger.error("parser.csv.error", path=file_path, error=str(exc))
            return []


class ExcelParser(BaseParser):
    """Parser for Excel (.xlsx, .xls) files."""

    supported_extensions = [".xlsx", ".xls"]
    supported_mimetypes = [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ]
    parser_name = "excel"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            chunks: list[ContentChunk] = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])

                if not rows:
                    continue

                # Convert to markdown table in batches
                header = rows[0]
                batch_size = 50
                for i in range(1, max(len(rows), 2), batch_size):
                    batch = rows[i:i + batch_size]
                    lines = ["| " + " | ".join(header) + " |"]
                    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in batch:
                        padded = row + [""] * (len(header) - len(row))
                        lines.append("| " + " | ".join(padded[:len(header)]) + " |")
                    chunks.append(ContentChunk(
                        content="\n".join(lines),
                        content_type="table",
                        section_title=f"Sheet: {sheet_name}",
                        metadata={
                            **metadata,
                            "parser": "excel",
                            "sheet": sheet_name,
                        },
                    ))

            wb.close()
            logger.info("parser.excel.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except Exception as exc:
            logger.error("parser.excel.error", path=file_path, error=str(exc))
            return []


class PowerPointParser(BaseParser):
    """Parser for PowerPoint (.pptx) files."""

    supported_extensions = [".pptx", ".ppt"]
    supported_mimetypes = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ]
    parser_name = "pptx"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        try:
            from pptx import Presentation  # type: ignore[import-untyped]

            prs = Presentation(file_path)
            chunks: list[ContentChunk] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                    if shape.has_table:
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(cells)
                        if rows:
                            header = rows[0]
                            lines = ["| " + " | ".join(header) + " |"]
                            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                            for r in rows[1:]:
                                lines.append("| " + " | ".join(r) + " |")
                            texts.append("\n".join(lines))

                if texts:
                    chunks.append(ContentChunk(
                        content="\n\n".join(texts),
                        content_type="text",
                        page_number=slide_num,
                        section_title=f"Slide {slide_num}",
                        metadata={**metadata, "parser": "pptx"},
                    ))

            logger.info("parser.pptx.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except ImportError:
            logger.error("parser.pptx.python_pptx_not_installed", path=file_path)
            return []
        except Exception as exc:
            logger.error("parser.pptx.error", path=file_path, error=str(exc))
            return []


class JSONParser(BaseParser):
    """Parser for JSON files."""

    supported_extensions = [".json"]
    supported_mimetypes = ["application/json"]
    parser_name = "json"

    async def parse(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        metadata = metadata or {}
        try:
            import json

            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            content = json.dumps(data, indent=2, ensure_ascii=False)
            md_parser = MarkdownParser()
            chunks = md_parser._chunk_by_size(
                f"```json\n{content}\n```",
                {**metadata, "parser": "json"},
            )
            logger.info("parser.json.parsed", path=file_path, chunks=len(chunks))
            return chunks

        except Exception as exc:
            logger.error("parser.json.error", path=file_path, error=str(exc))
            return []


# ── Parser Registry ──────────────────────────────────────────────────

class ParserRegistry:
    """Central registry that maps file types to parser implementations."""

    def __init__(self) -> None:
        self._parsers: list[BaseParser] = []
        self._register_defaults()

    def _register_defaults(self) -> None:
        """Register all built-in parsers."""
        self._parsers = [
            MarkdownParser(),
            PDFParser(),
            DocxParser(),
            HTMLParser(),
            CSVParser(),
            ExcelParser(),
            PowerPointParser(),
            JSONParser(),
        ]

    def register(self, parser: BaseParser) -> None:
        """Register a custom parser (inserted at highest priority)."""
        self._parsers.insert(0, parser)

    def get_parser(self, filename: str, mimetype: str = "") -> BaseParser | None:
        """Find the appropriate parser for the given file."""
        for parser in self._parsers:
            if parser.can_handle(filename, mimetype):
                return parser
        return None

    @property
    def supported_extensions(self) -> list[str]:
        """Return all supported file extensions."""
        exts: list[str] = []
        for parser in self._parsers:
            exts.extend(parser.supported_extensions)
        return sorted(set(exts))

    async def parse(
        self,
        file_path: str,
        filename: str = "",
        mimetype: str = "",
        metadata: dict[str, Any] | None = None,
    ) -> list[ContentChunk]:
        """Parse a file using the appropriate parser."""
        if not filename:
            filename = os.path.basename(file_path)

        parser = self.get_parser(filename, mimetype)
        if parser is None:
            logger.warning(
                "parser_registry.no_parser",
                filename=filename,
                mimetype=mimetype,
            )
            return []

        logger.info(
            "parser_registry.parsing",
            filename=filename,
            parser=parser.parser_name,
        )
        return await parser.parse(file_path, metadata)


# ── Singleton ────────────────────────────────────────────────────────

_registry: ParserRegistry | None = None


def get_parser_registry() -> ParserRegistry:
    """Return the singleton parser registry."""
    global _registry
    if _registry is None:
        _registry = ParserRegistry()
    return _registry
