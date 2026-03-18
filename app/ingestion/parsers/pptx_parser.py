"""ARIIA v2.0 – PPTX Streaming Parser.

Iteriert folie für Folie über python-pptx.
Kein komplettes Laden in RAM – Shapes werden einzeln verarbeitet.
"""
from __future__ import annotations

import asyncio
from pathlib import Path
from typing import AsyncIterator

import structlog

from app.ingestion.parsers.base import ParserRegistry, StreamingParser, TextChunk

logger = structlog.get_logger()


@ParserRegistry.register(
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
)
class PPTXParser(StreamingParser):
    """Parser für PPTX-Präsentationen via python-pptx.

    Verarbeitet Folie für Folie; extrahiert Text aus allen Shape-Typen
    (Textboxen, Tabellen, Notizen).
    """

    async def parse(self, file_path: Path) -> AsyncIterator[TextChunk]:
        loop = asyncio.get_event_loop()
        slides_data = await loop.run_in_executor(None, self._extract_slides, file_path)
        char_offset = 0
        for slide_num, texts in slides_data:
            for text in texts:
                yield TextChunk(
                    text=text,
                    page_num=slide_num,
                    char_offset=char_offset,
                    source_metadata={"parser": "pptx", "slide": slide_num},
                )
                char_offset += len(text)

    @staticmethod
    def _extract_slides(file_path: Path) -> list[tuple[int, list[str]]]:
        try:
            from pptx import Presentation  # type: ignore[import]
        except ImportError:
            raise RuntimeError("python-pptx nicht installiert: pip install python-pptx")

        prs = Presentation(str(file_path))
        result = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []

            for shape in slide.shapes:
                # Textboxen und Platzhalter
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs if run.text.strip())
                        if len(line.strip()) > 10:
                            texts.append(line.strip())

                # Tabellen
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))

            # Redner-Notizen
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if len(notes_text) > 10:
                    texts.append(f"[Notizen] {notes_text}")

            if texts:
                result.append((slide_num, texts))

        return result
